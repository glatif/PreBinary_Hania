"""
Document Processing Utilities for Quiz Generator

This module contains functions to extract text from various document formats
including PDF, DOCX, PPTX, and TXT files.
"""

import PyPDF2
import docx
import pptx
import io
import streamlit as st
from typing import List, Tuple, Optional


def extract_text_from_pdf(file) -> str:
    """
    Extracts text from a PDF file.
    
    Args:
        file: Uploaded file object from Streamlit
        
    Returns:
        str: Extracted text content
    """
    try:
        pdf_reader = PyPDF2.PdfReader(file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() or ""
        return text
    except Exception as e:
        return f"Error reading PDF: {e}"


def extract_text_from_docx(file) -> str:
    """
    Extracts text from a DOCX file.
    
    Args:
        file: Uploaded file object from Streamlit
        
    Returns:
        str: Extracted text content
    """
    try:
        doc = docx.Document(file)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        return f"Error reading DOCX: {e}"


def extract_text_from_pptx(file) -> str:
    """
    Extracts text from a PPTX file.
    
    Args:
        file: Uploaded file object from Streamlit
        
    Returns:
        str: Extracted text content
    """
    try:
        presentation = pptx.Presentation(file)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error reading PPTX: {e}"


def extract_text_from_txt(file) -> str:
    """
    Extracts text from a TXT file.
    
    Args:
        file: Uploaded file object from Streamlit
        
    Returns:
        str: Extracted text content
    """
    try:
        return file.getvalue().decode("utf-8")
    except Exception as e:
        return f"Error reading TXT: {e}"


def process_uploaded_files(uploaded_files) -> List[Tuple[str, str]]:
    """
    Process multiple uploaded files and extract text from each.
    
    Args:
        uploaded_files: List of uploaded file objects from Streamlit
        
    Returns:
        List[Tuple[str, str]]: List of tuples containing (filename, extracted_text)
    """
    extracted_texts = []
    
    if not uploaded_files:
        return extracted_texts
    
    for file in uploaded_files:
        try:
            file_extension = file.name.split('.')[-1].lower()
            text = ""
            
            if file_extension == "pdf":
                text = extract_text_from_pdf(file)
            elif file_extension == "docx":
                text = extract_text_from_docx(file)
            elif file_extension == "pptx":
                text = extract_text_from_pptx(file)
            elif file_extension == "txt":
                text = extract_text_from_txt(file)
            else:
                text = f"Unsupported file format: {file_extension}"
            
            if text and not text.startswith("Error"):
                extracted_texts.append((file.name, text))
            else:
                st.error(f"Failed to extract text from {file.name}: {text}")
                
        except Exception as e:
            st.error(f"Error processing {file.name}: {str(e)}")
    
    return extracted_texts


def combine_extracted_texts(extracted_texts: List[Tuple[str, str]]) -> str:
    """
    Combine all extracted texts into a single string for LLM processing.
    
    Args:
        extracted_texts: List of tuples containing (filename, text)
        
    Returns:
        str: Combined text content
    """
    if not extracted_texts:
        return ""
    
    combined_text = ""
    for filename, text in extracted_texts:
        combined_text += f"\n\n--- Content from {filename} ---\n\n"
        combined_text += text
        combined_text += "\n\n" + "="*50 + "\n\n"
    
    return combined_text


def validate_extracted_content(extracted_texts: List[Tuple[str, str]]) -> bool:
    """
    Validate that we have sufficient content to generate quiz questions.
    
    Args:
        extracted_texts: List of tuples containing (filename, text)
        
    Returns:
        bool: True if content is sufficient, False otherwise
    """
    if not extracted_texts:
        return False
    
    total_length = sum(len(text) for _, text in extracted_texts)
    return total_length >= 100  # Minimum 100 characters required
