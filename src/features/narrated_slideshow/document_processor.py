"""
Document Processor for Narrated Slideshow Feature

This module handles text extraction from PDF and PowerPoint files,
with validation for file size and page/slide limits.
"""

import os
from datetime import datetime
from typing import Dict, Any, Union
import streamlit as st

# Import existing PDF utilities
from src.utils.pdf_utils import extract_text_from_pdf

try:
    import pptx
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

def log_debug(phase: str, status: str, message: str) -> None:
    """Log debug messages with consistent formatting"""
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    print(f"[NARRATED_SLIDESHOW] [{timestamp}] [{phase}] [{status}] {message}")

def validate_file_limits(uploaded_file) -> Dict[str, Any]:
    """
    Validate uploaded file against size and page/slide limits
    
    Args:
        uploaded_file: Streamlit uploaded file object
        
    Returns:
        Dict with validation results
    """
    log_debug("VALIDATION", "INFO", f"Validating file: {uploaded_file.name}")
    
    file_type = uploaded_file.type.lower()
    file_size_mb = uploaded_file.size / (1024 * 1024)
    
    log_debug("VALIDATION", "INFO", f"File size: {file_size_mb:.2f} MB")
    
    # Basic file size check (100MB limit)
    if file_size_mb > 100:
        return {
            "valid": False,
            "message": f"File too large ({file_size_mb:.1f}MB). Maximum size is 100MB."
        }
    
    # Check file type
    if 'pdf' in file_type:
        return validate_pdf_limits(uploaded_file)
    elif 'presentation' in file_type or 'powerpoint' in file_type or file_type.endswith('.pptx') or file_type.endswith('.ppt'):
        return validate_ppt_limits(uploaded_file)
    else:
        return {
            "valid": False,
            "message": "Unsupported file type. Please upload a PDF or PowerPoint file."
        }

def validate_pdf_limits(uploaded_file) -> Dict[str, Any]:
    """Validate PDF file page limits"""
    try:
        import pypdf
        reader = pypdf.PdfReader(uploaded_file)
        page_count = len(reader.pages)
        
        log_debug("VALIDATION", "INFO", f"PDF has {page_count} pages")
        
        if page_count > 25:
            return {
                "valid": True,
                "message": f"PDF file ready for processing ({page_count} pages)",
                "warning": f"File has {page_count} pages. Only the first 25 pages will be processed."
            }
        else:
            return {
                "valid": True,
                "message": f"PDF file ready for processing ({page_count} pages)"
            }
            
    except Exception as e:
        log_debug("VALIDATION", "ERROR", f"Error validating PDF: {str(e)}")
        return {
            "valid": False,
            "message": f"Error reading PDF file: {str(e)}"
        }

def validate_ppt_limits(uploaded_file) -> Dict[str, Any]:
    """Validate PowerPoint file slide limits"""
    if not PPTX_AVAILABLE:
        return {
            "valid": False,
            "message": "PowerPoint processing not available. Please install python-pptx package."
        }
    
    try:
        # Save file temporarily to read with python-pptx
        temp_path = f"temp_{uploaded_file.name}"
        with open(temp_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        
        try:
            prs = Presentation(temp_path)
            slide_count = len(prs.slides)
            
            log_debug("VALIDATION", "INFO", f"PowerPoint has {slide_count} slides")
            
            if slide_count > 20:
                return {
                    "valid": True,
                    "message": f"PowerPoint file ready for processing ({slide_count} slides)",
                    "warning": f"File has {slide_count} slides. Only the first 20 slides will be processed."
                }
            else:
                return {
                    "valid": True,
                    "message": f"PowerPoint file ready for processing ({slide_count} slides)"
                }
        finally:
            # Clean up temp file
            if os.path.exists(temp_path):
                os.remove(temp_path)
                
    except Exception as e:
        log_debug("VALIDATION", "ERROR", f"Error validating PowerPoint: {str(e)}")
        return {
            "valid": False,
            "message": f"Error reading PowerPoint file: {str(e)}"
        }

def extract_pdf_content(uploaded_file) -> Dict[str, Any]:
    """
    Extract text content from PDF file page by page
    
    Args:
        uploaded_file: Streamlit uploaded file object
        
    Returns:
        Dict with extraction results
    """
    log_debug("EXTRACTION", "INFO", "Starting PDF text extraction")
    
    try:
        import pypdf
        
        # Reset file pointer to beginning
        uploaded_file.seek(0)
        
        reader = pypdf.PdfReader(uploaded_file)
        total_pages = len(reader.pages)
        
        # Limit to first 25 pages
        pages_to_process = min(total_pages, 25)
        log_debug("EXTRACTION", "INFO", f"Processing {pages_to_process} of {total_pages} pages")
        
        content = {}
        
        for page_num in range(pages_to_process):
            try:
                page = reader.pages[page_num]
                text = page.extract_text()
                
                # Clean up the text
                text = text.strip()
                content[page_num + 1] = text  # 1-indexed page numbers
                
                char_count = len(text)
                log_debug("EXTRACTION", "INFO", f"Page {page_num + 1}: {char_count} characters extracted")
                
                # Log a sample of the text for debugging
                if char_count > 0:
                    sample_text = text[:100] + "..." if len(text) > 100 else text
                    log_debug("EXTRACTION", "DEBUG", f"Page {page_num + 1} sample: {sample_text}")
                
            except Exception as e:
                log_debug("EXTRACTION", "WARNING", f"Error extracting page {page_num + 1}: {str(e)}")
                content[page_num + 1] = f"[Error extracting content from page {page_num + 1}]"
        
        log_debug("EXTRACTION", "SUCCESS", f"Successfully extracted content from {len(content)} pages")
        
        return {
            "success": True,
            "content": content,
            "total_pages": total_pages,
            "processed_pages": pages_to_process
        }
        
    except Exception as e:
        error_msg = f"Error extracting PDF content: {str(e)}"
        log_debug("EXTRACTION", "ERROR", error_msg)
        return {
            "success": False,
            "error": error_msg
        }

def extract_ppt_content(uploaded_file) -> Dict[str, Any]:
    """
    Extract text content from PowerPoint file slide by slide
    
    Args:
        uploaded_file: Streamlit uploaded file object
        
    Returns:
        Dict with extraction results
    """
    log_debug("EXTRACTION", "INFO", "Starting PowerPoint text extraction")
    
    if not PPTX_AVAILABLE:
        error_msg = "PowerPoint processing not available. Please install python-pptx package."
        log_debug("EXTRACTION", "ERROR", error_msg)
        return {
            "success": False,
            "error": error_msg
        }
    
    try:
        # Save file temporarily to read with python-pptx
        temp_path = f"temp_{uploaded_file.name}"
        with open(temp_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        
        try:
            prs = Presentation(temp_path)
            
            # Check if presentation has slides
            if not hasattr(prs, 'slides') or len(prs.slides) == 0:
                log_debug("EXTRACTION", "ERROR", "No slides found in presentation")
                return {
                    "success": False,
                    "error": "No slides found in presentation"
                }
            
            total_slides = len(prs.slides)
            
            # Limit to first 20 slides
            slides_to_process = min(total_slides, 20)
            log_debug("EXTRACTION", "INFO", f"Processing {slides_to_process} of {total_slides} slides")
            
            content = {}
            
            # Process slides one by one, following the working pattern from reference files
            for slide_num in range(slides_to_process):
                try:
                    slide = prs.slides[slide_num]
                    # Extract text using the exact working method from your streamlit_app_for_ppt.py
                    text_runs = []
                    
                    # Check if slide has shapes
                    if not hasattr(slide, 'shapes'):
                        log_debug("EXTRACTION", "WARNING", f"Slide {slide_num + 1} has no shapes")
                        content[slide_num + 1] = ""
                        continue
                    
                    for shape in slide.shapes:
                        if not hasattr(shape, 'has_text_frame') or not shape.has_text_frame:
                            continue
                        if not hasattr(shape, 'text_frame') or not shape.text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text_runs.append(run.text)
                    
                    # Combine all text from the slide
                    slide_content = "".join(text_runs).strip()
                    content[slide_num + 1] = slide_content  # 1-indexed slide numbers
                    
                    char_count = len(slide_content)
                    log_debug("EXTRACTION", "INFO", f"Slide {slide_num + 1}: {char_count} characters extracted")
                    
                    # Log a sample of the text for debugging
                    if char_count > 0:
                        sample_text = slide_content[:100] + "..." if len(slide_content) > 100 else slide_content
                        log_debug("EXTRACTION", "DEBUG", f"Slide {slide_num + 1} sample: {sample_text}")
                    
                except Exception as e:
                    log_debug("EXTRACTION", "WARNING", f"Error extracting slide {slide_num + 1}: {str(e)}")
                    content[slide_num + 1] = f"[Error extracting content from slide {slide_num + 1}]"
            
            log_debug("EXTRACTION", "SUCCESS", f"Successfully extracted content from {len(content)} slides")
            
            return {
                "success": True,
                "content": content,
                "total_slides": total_slides,
                "processed_slides": slides_to_process
            }
            
        finally:
            # Clean up temp file
            if os.path.exists(temp_path):
                os.remove(temp_path)
                
    except Exception as e:
        error_msg = f"Error extracting PowerPoint content: {str(e)}"
        log_debug("EXTRACTION", "ERROR", error_msg)
        return {
            "success": False,
            "error": error_msg
        }

def process_uploaded_file(uploaded_file) -> Dict[str, Any]:
    """
    Process uploaded file and extract content based on file type
    
    Args:
        uploaded_file: Streamlit uploaded file object
        
    Returns:
        Dict with processing results
    """
    log_debug("PROCESSING", "INFO", f"Processing file: {uploaded_file.name}")
    
    file_type = uploaded_file.type.lower()
    
    if 'pdf' in file_type:
        return extract_pdf_content(uploaded_file)
    elif 'presentation' in file_type or 'powerpoint' in file_type or file_type.endswith('.pptx') or file_type.endswith('.ppt'):
        return extract_ppt_content(uploaded_file)
    else:
        error_msg = f"Unsupported file type: {file_type}"
        log_debug("PROCESSING", "ERROR", error_msg)
        return {
            "success": False,
            "error": error_msg
        }
