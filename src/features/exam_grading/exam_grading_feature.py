# =============================================================================
# exam_grading_feature.py — Exam Grading Feature
# =============================================================================
# Provides the full Exam Grading UI and all supporting database operations.
#
# Feature overview:
#   Teachers set up an exam by providing questions and a grading rubric, upload
#   student submissions (individual PDFs or a ZIP archive), and trigger an LLM
#   grading pass. Results are displayed in real time and persisted to the
#   database so they can be reviewed later via the History tab.
#
# Database integration:
#   - Each graded student result is saved to exam_grading_results, linked to
#     the current assessment via assessment_id.
#   - The History tab queries only the sessions belonging to the current
#     assessment, keeping history scoped to the context the teacher is working in.
#   - The History tab Load Setup button restores the questions, rubric, and
#     max points from any previous session so the teacher can reuse the same
#     exam setup with new student submissions.
#
# Course/assessment context:
#   This feature is always rendered inside a specific assessment. The selected
#   course and assessment are read from session state using the tab-namespaced
#   keys set by app.py's navigation system:
#     st.session_state["exam_grading_selected_course"]     -> {"id": int, "name": str}
#     st.session_state["exam_grading_selected_assessment"] -> {"id": int, "title": str}
# =============================================================================

import os
import json
import uuid
import streamlit as st
import pandas as pd
import zipfile
import tempfile
from typing import List, Dict, Any
from db import get_connection
from auth import save_uploaded_file

from src.utils.llm_utils import MODELS, generate_llm_response, MODEL_PROVIDERS, strip_llm_json
from src.utils.pdf_utils import extract_text_from_pdf, save_uploaded_pdf
from src.features.exam_verification.exam_verification_feature import verify_student_identity
from src.features.proctoring.proctoring_feature import (
    render_proctor_monitor,
    get_proctor_summary_by_user_assessment,
    get_proctor_frames_by_user_assessment,
)

try:
    from docx import Document as DocxDocument
    _DOCX_AVAILABLE = True
except ImportError:
    _DOCX_AVAILABLE = False

try:
    from pptx import Presentation as PptxPresentation
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

SUPPORTED_SUBMISSION_TYPES = ["pdf", "doc", "docx", "txt", "pptx", "ppt"]


# =============================================================================
# DATABASE WRITE OPERATIONS
# =============================================================================

def save_exam_grading_result(
    grading_session_id: str,
    graded_by: int,
    assessment_id: int,
    student_name: str,
    student_id_parsed: str,
    questions_text: str,
    rubric: str,
    sub_rubric: str,
    score: float,
    max_points: int,
    feedback: str,
    detailed_explanation: str,
    model_name: str,
) -> None:
    """
    Persist one graded student result to the exam_grading_results table.

    One row is written per student submission. All rows produced by a single
    grading run share the same grading_session_id so the History tab can group
    and display them as a batch.

    assessment_id links the result to the assessment the teacher was working
    under when the grading run was triggered. This allows the History tab to
    filter by assessment rather than showing all sessions across all courses.

    The questions_text, rubric, and sub_rubric used for the session are stored
    redundantly per row to maintain a self-contained audit record for each
    individual student result.
    """
    model_provider = MODEL_PROVIDERS.get(model_name)

    conn = get_connection()
    cursor = conn.cursor()
    try:
        cursor.execute("""
            INSERT INTO exam_grading_results (
                grading_session_id,
                graded_by,
                assessment_id,
                student_name,
                student_id_parsed,
                questions_text,
                rubric,
                sub_rubric,
                score,
                max_points,
                feedback,
                detailed_explanation,
                model_provider,
                model_name
            )
            VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
        """, (
            grading_session_id,
            graded_by,
            assessment_id,
            student_name,
            student_id_parsed,
            questions_text,
            rubric,
            sub_rubric,
            score,
            max_points,
            feedback,
            detailed_explanation,
            model_provider,
            model_name,
        ))
        conn.commit()
    finally:
        cursor.close()
        conn.close()


# =============================================================================
# STUDENT SUBMISSION — "Submit My Exam"
# =============================================================================
# Students upload their own completed exam file directly, after passing the
# ID-card + live-selfie identity check. Files are written through the same
# save_uploaded_file()/files-table mechanism every other feature uses, with
# feature_name='exam_grading_submission' so the teacher's Student Submissions
# tab can find exactly these rows. This is required (rather than session
# state) because the student and the teacher are in separate browser
# sessions — st.session_state cannot bridge across users.

def save_student_exam_submission(
    student_id: int,
    assessment_id: int,
    course_id: int,
    file_bytes: bytes,
    original_name: str,
    course_name: str,
    assessment_name: str,
) -> None:
    """Persist a verified student's own exam submission file."""
    saved_name, saved_path = save_uploaded_file(
        file_bytes=file_bytes,
        original_name=original_name,
        course_name=course_name,
        assessment_name=assessment_name,
        course_id=course_id,
        feature_name="exam_grading_submission",
    )
    conn = get_connection()
    cursor = conn.cursor()
    try:
        cursor.execute(
            """
            INSERT INTO files (file_name, file_path, course_id, assessment_id, uploaded_by, feature_name)
            VALUES (%s, %s, %s, %s, %s, %s)
            """,
            (saved_name, saved_path, course_id, assessment_id, student_id, "exam_grading_submission"),
        )
        conn.commit()
    finally:
        cursor.close()
        conn.close()


def get_student_exam_submissions(assessment_id: int, student_id: int = None) -> List[Dict]:
    """
    Return exam_grading_submission file rows for an assessment.

    Filtered to one student's own files when student_id is given (the
    student's own "already submitted" list), or all students' files when
    omitted (the teacher's Student Submissions tab).
    """
    conn = get_connection()
    cursor = conn.cursor(dictionary=True)
    try:
        query = """
            SELECT f.id, f.file_name, f.file_path, f.uploaded_at, f.uploaded_by,
                   u.first_name, u.last_name, u.roll_no
            FROM files f
            JOIN users u ON u.id = f.uploaded_by
            WHERE f.assessment_id = %s AND f.feature_name = 'exam_grading_submission'
        """
        params = [assessment_id]
        if student_id is not None:
            query += " AND f.uploaded_by = %s"
            params.append(student_id)
        query += " ORDER BY f.uploaded_at DESC"
        cursor.execute(query, tuple(params))
        return cursor.fetchall()
    finally:
        cursor.close()
        conn.close()


def save_exam_setup(
    assessment_id: int,
    questions: str,
    rubric: str,
    sub_rubric: str,
    max_points: int,
    set_by: int,
) -> None:
    """
    Persist the canonical exam setup for an assessment so students can read
    the questions (rubric/sub_rubric are stored for the grading pipeline but
    are never shown to students). One row per assessment — upserted on every
    save, since a teacher revising the exam should replace the prior version
    rather than accumulate duplicates.
    """
    conn = get_connection()
    cursor = conn.cursor()
    try:
        cursor.execute(
            """
            INSERT INTO exam_setups (assessment_id, questions, rubric, sub_rubric, max_points, set_by)
            VALUES (%s, %s, %s, %s, %s, %s)
            ON DUPLICATE KEY UPDATE
                questions  = VALUES(questions),
                rubric     = VALUES(rubric),
                sub_rubric = VALUES(sub_rubric),
                max_points = VALUES(max_points),
                set_by     = VALUES(set_by)
            """,
            (assessment_id, questions, rubric, sub_rubric, max_points, set_by),
        )
        conn.commit()
    finally:
        cursor.close()
        conn.close()


def get_exam_setup(assessment_id: int) -> Dict:
    """Return the saved exam setup for an assessment, or None if not yet saved."""
    conn = get_connection()
    cursor = conn.cursor(dictionary=True)
    try:
        cursor.execute(
            "SELECT questions, rubric, sub_rubric, max_points, updated_at "
            "FROM exam_setups WHERE assessment_id = %s",
            (assessment_id,),
        )
        return cursor.fetchone()
    finally:
        cursor.close()
        conn.close()


def _render_student_exam_submission(
    course_id: int,
    course_name: str,
    assessment_id: int,
    assessment_title: str,
) -> None:
    """
    Student-facing "Submit My Exam" view, rendered instead of the full
    teacher workflow when exam_grading_ui() detects a student role.

    Exam questions are shown up front (no verification needed to read them —
    a student may want to prepare their answers before opening the camera).
    The upload form itself is gated by verify_student_identity() — it only
    appears once the student has shown their ID card and a live selfie that
    match their profile. One verification per assessment (cached in session
    state by the gate_key) covers re-submission attempts within the same
    session. The rubric/sub_rubric are intentionally never shown here — that
    is the grading key.

    Once verified, render_proctor_monitor() starts tab-switch/focus-loss
    monitoring and the one-click screen-share prompt for the rest of this
    render — the student has the questions in front of them and could switch
    away to search for or draft answers before uploading, so monitoring
    covers the whole window between verification and submission, not just
    the upload click itself.
    """
    user = st.session_state.user

    st.markdown("### 📤 Submit My Exam")

    if not assessment_id:
        st.warning("Select a course and assessment first.")
        return

    setup = get_exam_setup(assessment_id)
    if not setup or not setup.get("questions"):
        st.info("Your instructor has not published the exam questions yet. Check back later.")
        return

    with st.expander("📋 Exam Questions", expanded=True):
        st.caption(f"Maximum points: {setup.get('max_points', 100)}")
        st.markdown(setup["questions"])

    st.divider()

    if not verify_student_identity(user, gate_key=f"exam_grading_{assessment_id}"):
        return

    render_proctor_monitor(
        gate_key=f"exam_grading_{assessment_id}",
        user=user,
        quiz_id=None,
        assessment_id=assessment_id,
    )

    st.success("Identity verified. You may now upload your exam.")

    existing = get_student_exam_submissions(assessment_id, student_id=int(user["id"]))
    if existing:
        st.info(f"You have already submitted {len(existing)} file(s) for this assessment.")
        for row in existing:
            st.caption(f"`{row['file_name']}` — submitted {row['uploaded_at']}")

    uploaded_file = st.file_uploader(
        "Upload your completed exam (PDF, Word, DOC, TXT, or PowerPoint)",
        type=SUPPORTED_SUBMISSION_TYPES,
        key=f"student_exam_upload_{assessment_id}",
    )
    if uploaded_file and st.button(
        "Submit Exam", type="primary", key=f"submit_exam_btn_{assessment_id}"
    ):
        try:
            save_student_exam_submission(
                student_id=int(user["id"]),
                assessment_id=assessment_id,
                course_id=course_id,
                file_bytes=uploaded_file.read(),
                original_name=uploaded_file.name,
                course_name=course_name,
                assessment_name=assessment_title,
            )
            st.success("Exam submitted successfully.")
            st.rerun()
        except Exception as exc:
            st.error(f"Could not submit your exam: {exc}")


# =============================================================================
# DATABASE READ OPERATIONS
# =============================================================================

def get_exam_grading_sessions(user_id: int, assessment_id: int) -> List[Dict]:
    """
    Return all grading sessions run by the given user within the given assessment.

    Results are grouped by grading_session_id so that each row represents one
    batch run rather than one student. Filtering by assessment_id keeps the
    History tab scoped to the assessment the teacher currently has open.
    """
    conn = get_connection()
    cursor = conn.cursor(dictionary=True)
    try:
        cursor.execute("""
            SELECT
                grading_session_id,
                MAX(graded_at) AS graded_at,
                COUNT(*) AS student_count
            FROM exam_grading_results
            WHERE graded_by = %s
              AND assessment_id = %s
            GROUP BY grading_session_id
            ORDER BY graded_at DESC
        """, (user_id, assessment_id))
        return cursor.fetchall() or []
    finally:
        cursor.close()
        conn.close()


def get_exam_grading_session_results(grading_session_id: str, user_id: int) -> List[Dict]:
    """
    Return all individual student results belonging to a specific grading session.

    Filtered by both grading_session_id and graded_by so that a user can only
    retrieve results from their own sessions.
    """
    conn = get_connection()
    cursor = conn.cursor(dictionary=True)
    try:
        cursor.execute("""
            SELECT
                id,
                student_name,
                student_id_parsed,
                score,
                max_points,
                feedback,
                detailed_explanation,
                model_provider,
                model_name,
                graded_at
            FROM exam_grading_results
            WHERE grading_session_id = %s
              AND graded_by = %s
            ORDER BY graded_at DESC, student_name ASC
        """, (grading_session_id, user_id))
        return cursor.fetchall() or []
    finally:
        cursor.close()
        conn.close()


def get_exam_grading_session_setup(grading_session_id: str, user_id: int) -> Dict:
    """
    Return the exam setup data for a specific grading session.

    Fetches the questions_text, rubric, sub_rubric, and max_points from a
    single result row belonging to the session. All rows in a session share
    the same setup values, so only one row is needed.

    Returns an empty dict if no results are found for the session.
    """
    conn = get_connection()
    cursor = conn.cursor(dictionary=True)
    try:
        cursor.execute("""
            SELECT
                questions_text,
                rubric,
                sub_rubric,
                max_points
            FROM exam_grading_results
            WHERE grading_session_id = %s
              AND graded_by = %s
            LIMIT 1
        """, (grading_session_id, user_id))
        result = cursor.fetchone()
        return result or {}
    finally:
        cursor.close()
        conn.close()


def delete_exam_grading_session(grading_session_id: str, user_id: int) -> None:
    """
    Permanently delete all exam grading result rows for a single session.

    Filtered by both grading_session_id and graded_by so a teacher can only
    delete their own sessions. All student results for the session are removed
    in one DELETE — the session itself is not a separate table row.
    """
    conn = get_connection()
    cursor = conn.cursor()
    try:
        cursor.execute("""
            DELETE FROM exam_grading_results
            WHERE grading_session_id = %s
              AND graded_by = %s
        """, (grading_session_id, user_id))
        conn.commit()
    finally:
        cursor.close()
        conn.close()


def update_exam_grading_feedback(
    grading_session_id: str,
    student_name: str,
    feedback: str,
    detailed_explanation: str,
    graded_by: int,
) -> None:
    """
    Overwrite the feedback and detailed_explanation for a single student result.

    Identified by grading_session_id + student_name + graded_by so the update
    is scoped to the correct teacher and session.
    """
    conn = get_connection()
    cursor = conn.cursor()
    try:
        cursor.execute("""
            UPDATE exam_grading_results
               SET feedback = %s,
                   detailed_explanation = %s
             WHERE grading_session_id = %s
               AND student_name = %s
               AND graded_by = %s
        """, (feedback, detailed_explanation, grading_session_id, student_name, graded_by))
        conn.commit()
    finally:
        cursor.close()
        conn.close()


# =============================================================================
# LLM PROMPT BUILDER
# =============================================================================

def create_grading_prompt(
    question: str,
    rubric: str,
    sub_rubric: str,
    student_answer: str,
    max_points: int,
) -> str:
    """
    Build the LLM prompt used to grade a single student's submission.

    The prompt instructs the model to return a strict JSON object containing
    the student's name, ID, a one-line feedback summary, a detailed explanation,
    and a numerical score. The JSON template is embedded verbatim so the model
    has an unambiguous format to follow.

    The response is parsed in exam_grading_ui(). If the top-level parse fails,
    a fallback extraction attempts to locate the JSON object within any
    surrounding text the model may have added.
    """
    json_template = """{
  "student_name": "STUDENT_NAME",
  "student_id": "STUDENT_ID",
  "detailed_explanation": "thorough feedback referencing each rubric criterion and explaining what the student did well and what was missing",
  "feedback": "one concise sentence of constructive feedback for the student",
  "score": REPLACE_WITH_INTEGER
}"""
    return f"""You are an expert teacher grading an exam response.
Evaluate the student's answer strictly against the rubric and criteria below.

SCORING RULES:
- The score MUST be an integer between 0 and {max_points}.
- Award {max_points} (full marks) if the answer fully satisfies every criterion.
- Deduct marks only for specific missing or incorrect content, not for writing style.
- Do NOT apply a penalty if a criterion is not mentioned but is also not wrong.
- Every point deducted must be explained in detailed_explanation.

---
QUESTION:
{question}

GRADING RUBRIC:
{rubric}

DETAILED EVALUATION CRITERIA:
{sub_rubric}

---
STUDENT'S ANSWER:
{student_answer}

---
MAXIMUM POINTS: {max_points}

Respond with ONLY valid JSON matching this exact structure (replace the placeholder values):
{json_template}
"""


def format_questions_with_llm(raw_text: str, model: str) -> str:
    """
    Send raw extracted text to the LLM and ask it to identify, number, and
    return only the exam questions — stripping headers, footers, instructions,
    and any non-question content.

    Handles combined text from multiple files (separated by --- filename ---
    markers) including PPT slides where questions may appear as slide titles
    or bullet points rather than traditional Q&A formatting.
    """
    prompt = (
        "You are an exam creator extracting questions from uploaded documents.\n"
        "The text below may be separated by '--- filename ---' markers. "
        "You MUST process every single section — do not skip any.\n\n"
        "=== WHAT COUNTS AS A QUESTION ===\n"
        "Extract ANY of the following — all are valid exam questions:\n"
        "- Multiple-choice questions with options (A/B/C/D or a/b/c/d)\n"
        "- True / False statements\n"
        "- Short-answer questions (starting with: What, Who, When, Where, Which, How, Why)\n"
        "- Command-form questions (starting with: Describe, Explain, Define, List, "
        "Identify, Outline, Discuss, Compare, Contrast, Evaluate, Analyze, Justify, "
        "Calculate, Suggest, State, Summarise, Differentiate)\n"
        "- Application / scenario questions (starting with: Given, Consider, In the "
        "following, Suppose, Refer to, Based on, Using the)\n"
        "- Any numbered or lettered item that clearly expects a student to respond\n\n"
        "=== WHAT TO DO WITH EACH SECTION ===\n"
        "RULE A — Document already contains questions (any type above): "
        "copy every single question EXACTLY as written. "
        "Keep MCQ options on separate lines directly below their question.\n\n"
        "RULE B — Slide / lecture material (titles, bullets, concept lists with no "
        "question structure): convert every distinct point into a question.\n\n"
        "=== OUTPUT FORMAT ===\n"
        "- Number ALL questions sequentially (1. 2. 3. ...) across all sections.\n"
        "- Output ONLY the numbered questions — no file headers, no labels, "
        "no preamble, no explanations.\n\n"
        f"EXTRACTED TEXT:\n{raw_text}"
    )
    return generate_llm_response(prompt, model)


def _process_single_file_for_questions(filename: str, content: str, model: str) -> str:
    """
    Extract or generate questions from a single file's content.

    Two modes in one prompt:
    - Explicit question files (MCQ, Q&A, true/false): questions are copied as-is.
    - Material / slide files (titles, bullets, concepts): each point is converted
      into a question.

    Keeping files separate — rather than merging them before calling the LLM —
    prevents the model from treating slide content as context for MCQs and
    discarding it.
    """
    prompt = (
        f"You are extracting exam questions from the file: {filename}\n\n"
        "=== WHAT COUNTS AS A QUESTION ===\n"
        "The following are ALL valid exam question types — extract every one:\n"
        "- Multiple-choice questions with options (A/B/C/D)\n"
        "- True / False statements\n"
        "- Short-answer questions (What, Who, When, Where, Which, How, Why)\n"
        "- Command-form questions (Describe, Explain, Define, List, Identify, "
        "Outline, Discuss, Compare, Contrast, Evaluate, Analyze, Justify, "
        "Calculate, Suggest, State, Summarise, Differentiate)\n"
        "- Application / scenario questions (Given, Consider, In the following, "
        "Suppose, Refer to, Based on, Using the)\n"
        "- Any numbered or lettered item that expects a student response\n\n"
        "=== RULES ===\n"
        "RULE 1 — Content already has questions (any type above): "
        "copy every single one EXACTLY as written. "
        "Keep MCQ options on separate lines directly below their question.\n\n"
        "RULE 2 — Slide / lecture material (titles, bullets, concept lists): "
        "convert every distinct point into a question. "
        "Example — 'Types of RAM: SRAM, DRAM' → "
        "'What are the two main types of RAM? Describe each.'\n\n"
        "Output format: numbered list starting at 1 (1. 2. 3. ...)\n"
        "Output ONLY the questions — no preamble, no filename, no explanations.\n\n"
        f"CONTENT:\n{content}"
    )
    return generate_llm_response(prompt, model)


def _renumber_questions(text: str) -> str:
    """Renumber all question lines sequentially across a combined questions string."""
    import re
    lines = text.split("\n")
    result = []
    counter = 1
    for line in lines:
        stripped = line.strip()
        if re.match(r"^\d+[\.\)]\s", stripped):
            line = re.sub(r"^\d+[\.\)]\s*", f"{counter}. ", stripped)
            counter += 1
        result.append(line)
    return "\n".join(result)


def generate_questions_from_prompt(material: str, user_prompt: str, model: str) -> str:
    """
    Generate exam questions from study material according to a teacher-supplied prompt.

    The prompt describes the desired question types and count
    (e.g. "10 MCQ and 5 true/false"). The model returns only the questions,
    cleanly numbered, with no answer key or preamble.
    """
    # Strip file-separator markers so the LLM sees clean content, not metadata.
    import re
    clean_material = re.sub(r"^---\s*.+?\s*---\s*$", "", material, flags=re.MULTILINE).strip()

    prompt = (
        "You are an expert exam question writer.\n\n"
        "=== YOUR TASK ===\n"
        f"{user_prompt}\n\n"
        "=== STRICT OUTPUT RULES ===\n"
        "1. Write BRAND NEW questions — do NOT copy or repeat the study material.\n"
        "2. Every question must be answerable using information from the study material.\n"
        "3. Number every question sequentially: 1. 2. 3. ...\n"
        "4. For multiple-choice questions write the question first, then options on "
        "separate lines: A) ... B) ... C) ... D) ...\n"
        "5. For true/false questions write: True/False: <statement>\n"
        "6. Output ONLY the questions — no headings, no answer key, no explanations, "
        "no preamble, no 'Here are your questions:' intro.\n\n"
        "=== STUDY MATERIAL (read this, then write questions FROM it) ===\n"
        f"{clean_material}"
    )
    return generate_llm_response(prompt, model)


# =============================================================================
# MULTI-FORMAT TEXT EXTRACTION HELPERS
# =============================================================================

def _extract_text_from_docx(file_path: str) -> str:
    if not _DOCX_AVAILABLE:
        raise RuntimeError("python-docx is not installed.")
    doc = DocxDocument(file_path)
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_text_from_pptx(file_path: str) -> str:
    if not _PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is not installed.")
    prs = PptxPresentation(file_path)
    lines = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_lines = []
        for shape in slide.shapes:
            # Text frames (titles, text boxes, content placeholders)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_lines.append(text)
            # Tables — questions are often laid out in table cells
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            slide_lines.append(text)
        if slide_lines:
            lines.append(f"[Slide {slide_num}]")
            lines.extend(slide_lines)
    return "\n".join(lines)


def _extract_text_from_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _extract_text_from_doc(file_path: str) -> str:
    """Extract text from legacy .doc (Word 97-2003) binary files.

    python-docx handles some .doc files that are internally XML-based.
    For true OLE binary .doc files it fails, so we fall back to scanning
    the raw bytes for printable ASCII runs — imperfect but sufficient for
    the LLM to reconstruct exam questions from.
    """
    if _DOCX_AVAILABLE:
        try:
            doc = DocxDocument(file_path)
            text = "\n".join(p.text for p in doc.paragraphs)
            if text.strip():
                return text
        except Exception:
            pass

    import re
    with open(file_path, "rb") as f:
        raw = f.read()
    chunks = re.findall(rb"[\x20-\x7e\t]{4,}", raw)
    return "\n".join(c.decode("ascii", errors="replace") for c in chunks)


def extract_text_from_file(file_path: str) -> str:
    """Dispatch text extraction based on file extension."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    if ext == ".doc":
        return _extract_text_from_doc(file_path)
    if ext == ".docx":
        return _extract_text_from_docx(file_path)
    if ext in (".pptx", ".ppt"):
        return _extract_text_from_pptx(file_path)
    if ext == ".txt":
        return _extract_text_from_txt(file_path)
    raise ValueError(f"Unsupported file type: {ext}")


# =============================================================================
# ZIP PROCESSING HELPER
# =============================================================================

def process_zip_file(uploaded_zip) -> List[Dict]:
    """
    Extract and parse student submissions from a ZIP archive.

    The ZIP is extracted into a temporary directory. All supported file types
    (pdf, docx, txt, pptx, ppt) found anywhere inside the archive — including
    inside sub-folders — are processed. The expected filename convention is
    Name_ID.<ext>; if an underscore is present the portion before it becomes
    the student name and the portion after the first underscore becomes the
    student ID.

    The temporary directory and all its contents are deleted automatically when
    the context manager exits.

    Returns a list of submission dicts, each containing:
        student_name, student_id, filename, content (extracted text)
    """
    supported_exts = {f".{t}" for t in SUPPORTED_SUBMISSION_TYPES}
    submissions = []

    with tempfile.TemporaryDirectory() as tmpdirname:
        zip_path = os.path.join(tmpdirname, "submissions.zip")
        with open(zip_path, "wb") as f:
            f.write(uploaded_zip.getbuffer())

        with zipfile.ZipFile(zip_path, "r") as zip_ref:
            zip_ref.extractall(tmpdirname)

        # Walk all subdirectories so zips that contain student folders are handled.
        for dirpath, _dirs, filenames in os.walk(tmpdirname):
            for filename in filenames:
                ext = os.path.splitext(filename)[1].lower()
                if ext not in supported_exts:
                    continue
                file_path = os.path.join(dirpath, filename)
                try:
                    text = extract_text_from_file(file_path)

                    stem = os.path.splitext(filename)[0]
                    student_name = stem
                    student_id = ""
                    if "_" in stem:
                        parts = stem.split("_", 1)
                        student_name = parts[0]
                        student_id = parts[1]

                    submissions.append({
                        "student_name": student_name,
                        "student_id":   student_id,
                        "filename":     filename,
                        "content":      text,
                    })
                except Exception as e:
                    st.error(f"Error processing {filename}: {str(e)}")

    return submissions


# =============================================================================
# QUESTION FILE EXTRACTION HELPER
# =============================================================================

def extract_questions_from_file(uploaded_file) -> str:
    """
    Extract the text content from an uploaded question file (PDF, DOCX, TXT, PPTX/PPT).

    The file is written to a temporary path, text is extracted, and the
    temporary file is removed.

    Returns the extracted text as a string, or an empty string on failure.
    """
    try:
        os.makedirs("temp_uploads", exist_ok=True)
        file_path = os.path.join("temp_uploads", uploaded_file.name)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        text = extract_text_from_file(file_path)
        if os.path.exists(file_path):
            os.remove(file_path)
        return text
    except Exception as e:
        st.error(f"Error extracting questions from file: {str(e)}")
        return ""


# =============================================================================
# MAIN UI
# =============================================================================

@st.dialog("Delete Grading Session")
def _dialog_delete_grading_session(grading_session_id: str, user_id: int) -> None:
    """
    Confirmation modal for permanently deleting a grading session.

    All student results belonging to the session are removed from the database.
    This action cannot be undone.
    """
    st.warning("Are you sure you want to delete this grading session? All student results will be permanently removed.")
    col1, col2 = st.columns(2)
    if col1.button("Delete", type="primary", key="eg_dialog_confirm_delete"):
        delete_exam_grading_session(grading_session_id, user_id)
        st.toast("Grading session deleted.")
        st.rerun()
    if col2.button("Cancel", key="eg_dialog_cancel_delete"):
        st.rerun()


def exam_grading_ui() -> None:
    """
    Render the full Exam Grading feature UI.

    The feature is always opened from within a specific course assessment. The
    current course and assessment are read from tab-namespaced session state
    keys so that grading results and file uploads are linked to the correct
    assessment in the database.

    Students get a single cut-down view instead of the tabs below: an
    identity-verification-gated "Submit My Exam" upload (see
    _render_student_exam_submission()). Everything from here down is the
    teacher/admin workflow.

    Tab structure:
        Setup Exam          -- enter questions, rubric, and point total.
        Student Submissions -- upload individual PDFs or a ZIP archive, or
                                pull in files students submitted themselves.
        Grading Results     -- trigger LLM grading and view results.
        History             -- browse past grading sessions for this assessment.
    """
    st.subheader("Exam Grading System")

    # Read the course and assessment context set by app.py's navigation. These
    # values are used when persisting files and grading results to the database.
    _course      = st.session_state.get("exam_grading_selected_course", {}) or {}
    _assessment  = st.session_state.get("exam_grading_selected_assessment", {}) or {}
    course_id        = _course.get("id")
    course_name      = _course.get("name", "")
    assessment_id    = _assessment.get("id")
    assessment_title = _assessment.get("title", "")

    if st.session_state.get("user", {}).get("role") == "student":
        _render_student_exam_submission(course_id, course_name, assessment_id, assessment_title)
        return

    # Session state keys for this feature. Generic names match the original
    # session-based implementation so no downstream references need updating.
    if "questions" not in st.session_state:
        st.session_state.questions = ""
    if "rubric" not in st.session_state:
        st.session_state.rubric = ""
    if "sub_rubric" not in st.session_state:
        st.session_state.sub_rubric = ""
    if "max_points" not in st.session_state:
        st.session_state.max_points = 100
    # Tracks which input method is selected in the Setup Exam tab. Set to
    # "Enter questions manually" when setup is loaded from history so that
    # the loaded questions text is immediately visible in the text area.
    if "questions_input_method" not in st.session_state:
        st.session_state.questions_input_method = "Upload a file with questions"
    if "eg_raw_extracted" not in st.session_state:
        st.session_state["eg_raw_extracted"] = ""
    if "eg_zip_file_texts" not in st.session_state:
        st.session_state["eg_zip_file_texts"] = []
    if "submissions" not in st.session_state:
        st.session_state.submissions = []
    if "graded_results" not in st.session_state:
        st.session_state.graded_results = []

    tab1, tab2, tab3, tab4 = st.tabs([
        "📝 Setup Exam",
        "📄 Student Submissions",
        "📊 Grading Results",
        "🕘 History",
    ])

    # -------------------------------------------------------------------------
    # TAB 1 — Setup Exam
    # -------------------------------------------------------------------------
    with tab1:
        with st.expander("How this works", expanded=not st.session_state.questions):
            st.markdown(
                "**Step 1 — Setup Exam (this tab)**  \n"
                "Enter the exam questions and grading rubric. You can upload a single "
                "question paper (PDF, Word, DOC, TXT, or PowerPoint), a **ZIP archive** "
                "containing multiple question documents across sub-folders, or type "
                "them manually. After extracting, use **AI Format Questions** to clean "
                "up the raw text into numbered questions.  \n\n"
                "**Step 2 — Student Submissions**  \n"
                "Upload each student's answer file (individual files or a ZIP). "
                "Name files as `StudentName_ID.pdf` so the system can identify them.  \n\n"
                "**Step 3 — Grading Results**  \n"
                "Click **Grade All Submissions** and the AI will score every student "
                "against your rubric.  \n\n"
                "**Step 4 — History**  \n"
                "Revisit past grading sessions and download results as CSV."
            )

        # The input method is stored in session state so that loading a setup
        # from history can switch the mode to "Enter questions manually",
        # making the loaded questions text immediately visible.
        _METHOD_OPTS = [
            "Upload a file with questions",
            "Upload a ZIP file (supports sub-folders)",
            "Enter questions manually",
        ]
        input_method = st.radio(
            "How would you like to input the questions?",
            _METHOD_OPTS,
            index=_METHOD_OPTS.index(
                st.session_state.questions_input_method
                if st.session_state.questions_input_method in _METHOD_OPTS
                else _METHOD_OPTS[0]
            ),
            key="eg_input_method_radio",
        )
        st.session_state.questions_input_method = input_method

        if input_method == "Upload a file with questions":
            uploaded_pdf = st.file_uploader(
                "Upload the exam question paper (PDF, Word, DOC, TXT, or PowerPoint)",
                type=SUPPORTED_SUBMISSION_TYPES,
            )
            if uploaded_pdf:
                if st.button("Extract Text"):
                    with st.spinner("Extracting text from file..."):
                        raw_text = extract_questions_from_file(uploaded_pdf)
                        st.session_state["eg_raw_extracted"] = raw_text
                        st.success("Text extracted — choose how to generate questions below.")

        elif input_method == "Upload a ZIP file (supports sub-folders)":
            uploaded_zip = st.file_uploader(
                "Upload a ZIP containing question documents "
                "(PDF, Word, DOC, TXT, PowerPoint — sub-folders are supported)",
                type="zip",
            )
            if uploaded_zip:
                if st.button("Extract Text from ZIP"):
                    with st.spinner("Extracting text from all documents in ZIP..."):
                        try:
                            file_texts = process_zip_file(uploaded_zip)
                            if file_texts:
                                combined = "\n\n".join(
                                    f"--- {s['filename']} ---\n{s['content']}"
                                    for s in file_texts
                                )
                                st.session_state["eg_raw_extracted"] = combined
                                st.session_state["eg_zip_file_texts"] = file_texts
                                st.success(
                                    f"Extracted text from {len(file_texts)} file(s) — "
                                    "check the preview below, then choose how to generate questions."
                                )
                            else:
                                st.warning("No supported files found in the ZIP.")
                        except Exception as e:
                            st.error(f"Error processing ZIP: {str(e)}")

                if st.session_state.get("eg_zip_file_texts"):
                    with st.expander("Extracted content per file (verify before generating)", expanded=False):
                        for s in st.session_state["eg_zip_file_texts"]:
                            st.markdown(f"**{s['filename']}**")
                            preview = s["content"].strip()
                            if not preview:
                                st.warning("No text extracted from this file.")
                            else:
                                st.text(preview[:800] + ("…" if len(preview) > 800 else ""))
                            st.divider()

        else:
            st.session_state.questions = st.text_area(
                "Enter all exam questions:",
                value=st.session_state.questions,
                height=200,
            )

        # ------------------------------------------------------------------
        # Question generation — shown after material has been extracted from
        # an uploaded file or ZIP (not shown for manual entry).
        # ------------------------------------------------------------------
        if st.session_state.get("eg_raw_extracted") and input_method != "Enter questions manually":
            st.divider()
            st.markdown("#### Create exam questions from extracted material")

            gen_mode = st.radio(
                "Question generation mode",
                [
                    "Format existing questions from material",
                    "Generate questions using a custom prompt",
                ],
                key="eg_gen_mode_radio",
                label_visibility="collapsed",
            )

            _model_keys  = list(MODELS.keys())
            _saved_model = st.session_state.get("exam_grading_selected_model", _model_keys[0])
            if _saved_model not in _model_keys:
                _saved_model = _model_keys[0]
            _model = MODELS[_saved_model]

            if gen_mode == "Format existing questions from material":
                st.caption(
                    "The AI will scan the extracted text, identify every question, "
                    "and return them cleanly numbered — removing headers, footers, and instructions."
                )
                if st.button("AI Format Questions", key="eg_fmt_btn"):
                    file_texts = st.session_state.get("eg_zip_file_texts", [])
                    if file_texts:
                        # Process each file independently so slide/material files
                        # are not discarded when merged with MCQ files.
                        all_sections = []
                        prog = st.progress(0)
                        for i, file_info in enumerate(file_texts):
                            content = file_info["content"].strip()
                            if not content:
                                continue
                            with st.spinner(f"Processing {file_info['filename']}…"):
                                result = _process_single_file_for_questions(
                                    file_info["filename"], content, _model
                                )
                                result = str(result).strip() if result is not None else ""
                                if result and result not in ("{}", "{ }"):
                                    all_sections.append(result)
                            prog.progress((i + 1) / len(file_texts))
                        prog.empty()
                        if all_sections:
                            formatted = _renumber_questions("\n\n".join(all_sections))
                            st.session_state.questions = formatted
                            st.success("Questions formatted by AI!")
                        else:
                            st.error(
                                "No questions could be extracted from any file. "
                                "Check the extraction preview and try again."
                            )
                    else:
                        with st.spinner("Asking AI to identify and format questions..."):
                            formatted = format_questions_with_llm(
                                st.session_state["eg_raw_extracted"], _model
                            )
                            formatted = str(formatted).strip() if formatted is not None else ""
                            if not formatted or formatted in ("{}", "{ }"):
                                st.error(
                                    "The AI could not extract questions from this material. "
                                    "Check that the document contains readable text, then try again."
                                )
                            else:
                                st.session_state.questions = formatted
                                st.success("Questions formatted by AI!")

            else:
                st.caption(
                    "Describe the questions you want the AI to create from the material. "
                    "For example: *Create 10 multiple choice questions and 5 true/false questions.*"
                )
                gen_prompt = st.text_area(
                    "Generation prompt",
                    placeholder=(
                        "e.g. Create 10 multiple choice questions (A/B/C/D) and "
                        "5 true/false questions covering the key concepts in this material."
                    ),
                    height=90,
                    key="eg_gen_prompt_input",
                    label_visibility="collapsed",
                )
                if st.button("Generate Questions", key="eg_gen_btn"):
                    if not gen_prompt.strip():
                        st.warning("Please enter a generation prompt first.")
                    else:
                        with st.spinner("Generating questions from material..."):
                            generated = generate_questions_from_prompt(
                                st.session_state["eg_raw_extracted"],
                                gen_prompt.strip(),
                                _model,
                            )
                            generated = str(generated).strip() if generated is not None else ""
                            if not generated or generated in ("{}", "{ }"):
                                st.error(
                                    "The AI could not generate questions. "
                                    "Try rephrasing your prompt or check that the material extracted correctly."
                                )
                            else:
                                st.session_state.questions = generated
                                st.success("Questions generated!")

        if st.session_state.questions:
            with st.expander("Review Questions", expanded=True):
                st.markdown(str(st.session_state.questions))

        st.session_state.rubric = st.text_area(
            "Enter the general grading rubric:",
            value=st.session_state.rubric,
            height=150,
        )

        st.session_state.sub_rubric = st.text_area(
            "Enter detailed evaluation criteria (sub-rubric):",
            value=st.session_state.sub_rubric,
            height=150,
            help="Specify detailed criteria for each question (e.g., what to specifically look for in each response)",
        )

        st.session_state.max_points = st.number_input(
            "Maximum points for this exam:",
            min_value=1,
            value=st.session_state.max_points,
        )

        setup_col1, setup_col2 = st.columns(2)
        with setup_col1:
            if st.session_state.questions and st.button(
                "💾 Save Exam Setup (visible to students)", type="primary"
            ):
                save_exam_setup(
                    assessment_id=assessment_id,
                    questions=st.session_state.questions,
                    rubric=st.session_state.rubric,
                    sub_rubric=st.session_state.sub_rubric,
                    max_points=int(st.session_state.max_points),
                    set_by=int(st.session_state["user"]["id"]),
                )
                st.success(
                    "Exam setup saved — students can now see the questions in "
                    "their Submit My Exam view."
                )
        with setup_col2:
            if st.session_state.questions and st.button("Clear Exam Setup"):
                st.session_state.questions = ""
                st.session_state.rubric = ""
                st.session_state.sub_rubric = ""
                st.session_state.max_points = 100
                st.session_state["eg_raw_extracted"] = ""
                st.success("Exam setup cleared — upload a new question file or enter questions manually above.")

    # -------------------------------------------------------------------------
    # TAB 2 — Student Submissions
    # -------------------------------------------------------------------------
    with tab2:
        st.write("Upload student submissions")

        # ---- Files students submitted themselves via "Submit My Exam" ----
        # These come from a different browser session per student, so they
        # are persisted to the files table rather than session state — pull
        # them in here rather than requiring the teacher to re-collect and
        # re-upload files students already submitted directly.
        if assessment_id:
            student_files = get_student_exam_submissions(assessment_id)
            if not student_files:
                st.info(
                    "No students have submitted through \"Submit My Exam\" for "
                    "this assessment yet."
                )
            else:
                with st.expander(f"📥 Student-Submitted Files ({len(student_files)})", expanded=True):
                    st.caption(
                        "Uploaded directly by students after passing identity verification. "
                        "Load them into the grading queue below."
                    )
                    for row in student_files:
                        name = f"{row.get('first_name', '')} {row.get('last_name', '')}".strip() or "Unknown"
                        roll_suffix = f" (Roll No: {row['roll_no']})" if row.get("roll_no") else ""
                        st.write(f"**{name}**{roll_suffix} — `{row['file_name']}`")

                        # Tab-switch/focus-loss and screen-share summary recorded
                        # between this student's identity verification and their
                        # upload, aggregated across all their proctoring sessions
                        # for this assessment (see get_proctor_summary_by_user_assessment).
                        proctor = get_proctor_summary_by_user_assessment(row["uploaded_by"], assessment_id)
                        share_label = {
                            "granted": "✅ granted",
                            "denied":  "❌ denied",
                            None:      "— not recorded",
                        }[proctor["screen_share"]]
                        violation_count = proctor["violation_count"]
                        violation_icon  = "🔴" if violation_count else "🟢"
                        st.caption(
                            f"{violation_icon} {violation_count} tab-switch/focus warning(s) — "
                            f"Screen share: {share_label}"
                        )

                        # Screen-share snapshots captured between verification
                        # and upload, downscaled JPEGs taken every
                        # CAPTURE_INTERVAL_MS — see proctoring_feature.py.
                        frames = get_proctor_frames_by_user_assessment(row["uploaded_by"], assessment_id)
                        if frames:
                            with st.expander(f"📷 Screen Capture Frames ({len(frames)})", expanded=False):
                                frame_cols = st.columns(4)
                                for i, frame in enumerate(frames):
                                    with frame_cols[i % 4]:
                                        st.image(frame["file_path"], caption=str(frame["captured_at"]))

                    if st.button("Load Student-Submitted Files into Grading Queue", key="load_student_files_btn"):
                        with st.spinner("Processing student-submitted files..."):
                            loaded = []
                            for row in student_files:
                                try:
                                    content = extract_text_from_file(row["file_path"])
                                    name = (
                                        f"{row.get('first_name', '')} {row.get('last_name', '')}".strip()
                                        or row["file_name"]
                                    )
                                    loaded.append({
                                        "student_name": name,
                                        "student_id":   row.get("roll_no") or "",
                                        "filename":     row["file_name"],
                                        "content":      content,
                                    })
                                except Exception as e:
                                    st.error(f"Error processing {row['file_name']}: {str(e)}")
                            st.session_state.submissions = loaded
                            st.success(f"Loaded {len(loaded)} student-submitted file(s) into the grading queue.")
                st.divider()

        if not st.session_state.questions or not st.session_state.rubric:
            st.warning("Please set up questions and rubric in the Setup Exam tab first.")
        else:
            upload_option = st.radio(
                "How would you like to upload student submissions?",
                [
                    "Upload individual files",
                    "Upload a ZIP file (supports sub-folders)",
                ],
            )

            if upload_option == "Upload individual files":
                uploaded_files = st.file_uploader(
                    "Upload student submissions (PDF, Word, TXT, or PowerPoint)",
                    type=SUPPORTED_SUBMISSION_TYPES,
                    accept_multiple_files=True,
                )

                if uploaded_files:
                    if st.button("Process Individual Submissions"):
                        with st.spinner("Processing submissions..."):
                            submissions = []
                            for uploaded_file in uploaded_files:
                                try:
                                    os.makedirs("temp_uploads", exist_ok=True)
                                    file_path = os.path.join("temp_uploads", uploaded_file.name)
                                    with open(file_path, "wb") as f:
                                        f.write(uploaded_file.getbuffer())

                                    text = extract_text_from_file(file_path)

                                    stem = os.path.splitext(uploaded_file.name)[0]
                                    student_name = stem
                                    student_id = ""
                                    if "_" in stem:
                                        parts = stem.split("_", 1)
                                        student_name = parts[0]
                                        student_id = parts[1]

                                    submissions.append({
                                        "student_name": student_name,
                                        "student_id":   student_id,
                                        "filename":     uploaded_file.name,
                                        "content":      text,
                                    })

                                    if os.path.exists(file_path):
                                        os.remove(file_path)

                                except Exception as e:
                                    st.error(f"Error processing {uploaded_file.name}: {str(e)}")

                            st.session_state.submissions = submissions
                            st.success(f"Successfully processed {len(submissions)} submissions")

            else:
                uploaded_zip = st.file_uploader(
                    "Upload a ZIP file containing student submissions "
                    "(PDF, Word, TXT, PowerPoint — sub-folders are supported)",
                    type="zip",
                )

                if uploaded_zip:
                    if st.button("Process ZIP Submissions"):
                        with st.spinner("Extracting and processing PDFs from ZIP file..."):
                            try:
                                submissions = process_zip_file(uploaded_zip)
                                st.session_state.submissions = submissions
                                st.success(
                                    f"Successfully processed {len(submissions)} submissions from ZIP file"
                                )
                            except Exception as e:
                                st.error(f"Error processing ZIP file: {str(e)}")

            if st.session_state.submissions:
                st.write("### Processed Submissions")
                for i, submission in enumerate(st.session_state.submissions):
                    student_name    = submission["student_name"]
                    student_id      = submission.get("student_id", "")
                    student_id_text = f"(ID: {student_id})" if student_id else ""
                    with st.expander(f"Submission: {student_name} {student_id_text}"):
                        st.write("**Content Preview:**")
                        st.text(
                            submission["content"][:500] + "..."
                            if len(submission["content"]) > 500
                            else submission["content"]
                        )

    # -------------------------------------------------------------------------
    # TAB 3 — Grading Results
    # -------------------------------------------------------------------------
    with tab3:
        st.write("Grade student submissions")

        if not st.session_state.questions or not st.session_state.rubric:
            st.warning("Please set up questions and rubric in the Setup Exam tab first.")
        elif not st.session_state.submissions:
            st.warning("Please upload and process student submissions in the Student Submissions tab.")
        else:
            # Model selection — honours the user's saved preference loaded at
            # login by _load_model_preferences() into exam_grading_selected_model.
            # A separate widget key (exam_grading_model_selectbox) prevents
            # Streamlit from overwriting the preference on first render.
            if "exam_grading_selected_model" not in st.session_state:
                st.session_state.exam_grading_selected_model = list(MODELS.keys())[0]
            model_keys  = list(MODELS.keys())
            saved_model = st.session_state.get("exam_grading_selected_model", model_keys[0])
            if saved_model not in model_keys:
                saved_model = model_keys[0]
            selected_model_key = st.selectbox(
                "Select the model to use for grading:",
                model_keys,
                index=model_keys.index(saved_model),
                key="exam_grading_model_selectbox",
            )
            st.session_state.exam_grading_selected_model = selected_model_key
            selected_model = MODELS[selected_model_key]

            if selected_model == "llama-3.3-70b-groq" and not st.session_state.get("groq_api_key"):
                st.warning("⚠️ Groq API key is required. Please add your API key in your profile settings.")
            if selected_model == "gemini-2.5-flash" and not st.session_state.get("gemini_api_key"):
                st.warning("⚠️ Gemini API key is required. Please add your API key in your profile settings.")

            if st.button("Grade All Submissions"):
                grading_session_id = str(uuid.uuid4())
                graded_by          = st.session_state["user"]["id"]
                graded_results     = []

                progress_bar  = st.progress(0)
                status_text   = st.empty()
                results_table = st.empty()

                total_steps     = len(st.session_state.submissions)
                completed_steps = 0

                for submission in st.session_state.submissions:
                    student_name = submission["student_name"]
                    student_id   = submission.get("student_id", "")
                    content      = submission["content"]

                    status_text.text(f"Grading submission from {student_name}...")

                    prompt = create_grading_prompt(
                        question=st.session_state.questions,
                        rubric=st.session_state.rubric,
                        sub_rubric=st.session_state.sub_rubric,
                        max_points=st.session_state.max_points,
                        student_answer=content,
                    )

                    try:
                        response = generate_llm_response(prompt, selected_model, force_json=True)

                        # strip_llm_json extracts the outermost JSON object via
                        # brace-matching, handling bare JSON, fenced JSON, and
                        # responses with surrounding prose or trailing content.
                        # If parsing still fails, fall back to a default result
                        # so the grading run is not interrupted.
                        try:
                            result_json = json.loads(strip_llm_json(response))
                            # Coerce score: the LLM may return a string ("18"),
                            # a float (18.0), null, or the literal placeholder text.
                            raw_score = result_json.get("score")
                            try:
                                coerced = int(float(str(raw_score)))
                                # Clamp to valid range
                                result_json["score"] = max(0, min(coerced, st.session_state.max_points))
                            except (TypeError, ValueError):
                                result_json["score"] = 0
                        except Exception:
                            result_json = {
                                "student_name":         student_name,
                                "student_id":           student_id,
                                "score":                0,
                                "max_points":           st.session_state.max_points,
                                "feedback":             "Error parsing response",
                                "detailed_explanation": response,
                            }

                        # Ensure all expected fields are present and non-empty.
                        # Use explicit falsy checks instead of setdefault so that
                        # empty-string values returned by the LLM are also replaced.
                        if not result_json.get("score") and result_json.get("score") != 0:
                            result_json["score"] = 0
                        if not result_json.get("max_points"):
                            result_json["max_points"] = st.session_state.max_points
                        if not result_json.get("feedback", "").strip():
                            result_json["feedback"] = "No feedback provided"
                        if not result_json.get("detailed_explanation", "").strip():
                            result_json["detailed_explanation"] = response
                        if not result_json.get("student_name", "").strip():
                            result_json["student_name"] = student_name
                        if not result_json.get("student_id", "").strip():
                            result_json["student_id"] = student_id

                        # Always use the filename-parsed name and ID as the
                        # authoritative values; do not allow the model to
                        # substitute them from the prompt template placeholders.
                        result_json["student_name"]       = student_name
                        result_json["student_id"]         = student_id
                        result_json["grading_session_id"] = grading_session_id

                        save_exam_grading_result(
                            grading_session_id=grading_session_id,
                            graded_by=graded_by,
                            assessment_id=assessment_id,
                            student_name=student_name,
                            student_id_parsed=student_id,
                            questions_text=st.session_state.questions,
                            rubric=st.session_state.rubric,
                            sub_rubric=st.session_state.sub_rubric,
                            score=result_json["score"],
                            max_points=result_json["max_points"],
                            feedback=result_json["feedback"],
                            detailed_explanation=result_json["detailed_explanation"],
                            model_name=selected_model,
                        )

                        graded_results.append(result_json)

                    except Exception as e:
                        st.error(f"Error grading {student_name}'s submission: {str(e)}")
                        graded_results.append({
                            "student_name":         student_name,
                            "student_id":           student_id,
                            "score":                0,
                            "max_points":           st.session_state.max_points,
                            "feedback":             f"Error: {str(e)}",
                            "detailed_explanation": f"Failed to process submission: {str(e)}",
                        })

                    completed_steps += 1
                    progress_bar.progress(completed_steps / total_steps)

                    # Update the live results table after each student is graded.
                    if graded_results:
                        display_data = []
                        for r in graded_results:
                            sid_display = f"(ID: {r['student_id']})" if r.get("student_id") else ""
                            display_data.append({
                                "Student":              f"{r['student_name']} {sid_display}",
                                "Score":                f"{r['score']}/{r['max_points']}",
                                "Percentage":           f"{(float(r['score']) / float(r['max_points']) * 100):.1f}%",
                                "Feedback":             r.get("feedback", "N/A"),
                                "Detailed Explanation": r.get("detailed_explanation", "N/A"),
                            })
                        results_table.table(pd.DataFrame(display_data))

                st.session_state.graded_results = graded_results
                status_text.text("Grading completed!")
                progress_bar.progress(1.0)
                st.success("All submissions graded successfully!")

            # Display the final summary and per-student detail after grading.
            if st.session_state.graded_results:
                st.write("### Grading Results")

                summary_data = []
                for r in st.session_state.graded_results:
                    student_name    = r["student_name"]
                    student_id      = r.get("student_id", "")
                    student_id_text = f"(ID: {student_id})" if student_id else ""
                    score           = r["score"]
                    max_points      = r["max_points"]
                    percentage      = float(score) / float(max_points) * 100
                    summary_data.append({
                        "Student":    f"{student_name} {student_id_text}",
                        "Score":      f"{score}/{max_points}",
                        "Percentage": f"{percentage:.1f}%",
                        "Feedback":   r.get("feedback", "N/A"),
                    })

                if summary_data:
                    st.table(pd.DataFrame(summary_data))

                _current_user_id = st.session_state["user"]["id"]
                for i, result in enumerate(st.session_state.graded_results):
                    sid_display     = f"(ID: {result['student_id']})" if result.get("student_id") else ""
                    student_display = f"{result['student_name']} {sid_display}"
                    score_display   = f"{result['score']}/{result['max_points']}"
                    percentage      = float(result["score"]) / float(result["max_points"]) * 100

                    with st.expander(
                        f"Detailed Feedback: {student_display} — {score_display} ({percentage:.1f}%)"
                    ):
                        edited_feedback = st.text_area(
                            "Feedback",
                            value=result.get("feedback", ""),
                            height=80,
                            key=f"eg_edit_feedback_{i}",
                        )
                        edited_explanation = st.text_area(
                            "Detailed Explanation",
                            value=result.get("detailed_explanation", ""),
                            height=180,
                            key=f"eg_edit_explanation_{i}",
                        )
                        if st.button("💾 Save Changes", key=f"eg_save_feedback_{i}"):
                            st.session_state.graded_results[i]["feedback"]            = edited_feedback
                            st.session_state.graded_results[i]["detailed_explanation"] = edited_explanation
                            session_id_for_update = result.get("grading_session_id", "")
                            if session_id_for_update:
                                update_exam_grading_feedback(
                                    grading_session_id=session_id_for_update,
                                    student_name=result["student_name"],
                                    feedback=edited_feedback,
                                    detailed_explanation=edited_explanation,
                                    graded_by=_current_user_id,
                                )
                            st.toast(f"Feedback saved for {result['student_name']}.")

                # CSV export of the current session's results.
                if st.button("Download Results as CSV"):
                    try:
                        export_data = []
                        for r in st.session_state.graded_results:
                            export_data.append({
                                "Student Name":      r["student_name"],
                                "Student ID":        r.get("student_id", ""),
                                "Score":             r["score"],
                                "Max Points":        r["max_points"],
                                "Percentage":        float(r["score"]) / float(r["max_points"]) * 100,
                                "Feedback":          r.get("feedback", "N/A"),
                                "Detailed Feedback": r.get("detailed_explanation", "N/A"),
                            })
                        csv = pd.DataFrame(export_data).to_csv(index=False)
                        st.download_button(
                            label="Download CSV",
                            data=csv,
                            file_name="grading_results.csv",
                            mime="text/csv",
                        )
                    except Exception as e:
                        st.error(f"Error exporting results: {str(e)}")

    # -------------------------------------------------------------------------
    # TAB 4 — History
    # -------------------------------------------------------------------------
    with tab4:
        st.subheader("Grading History")

        user_id = st.session_state["user"]["id"]

        # History is scoped to the current assessment so teachers only see
        # sessions they ran while working in this specific assessment.
        sessions = get_exam_grading_sessions(user_id, assessment_id)

        if not sessions:
            st.info("No grading history found for this assessment.")
        else:
            for session in sessions:
                session_id    = session["grading_session_id"]
                graded_at     = session["graded_at"]
                student_count = session["student_count"]

                with st.expander(
                    f"{graded_at} — {student_count} student(s)", expanded=False
                ):
                    st.write(f"**Session ID:** {session_id}")
                    st.write(f"**Graded At:** {graded_at}")
                    st.write(f"**Students Graded:** {student_count}")

                    session_results = get_exam_grading_session_results(session_id, user_id)

                    if not session_results:
                        st.info("No student results found for this session.")
                    else:
                        for result in session_results:
                            student_label = (
                                f"{result.get('student_name', 'Unknown Student')} "
                                f"({result.get('score', 0)}/{result.get('max_points', 0)})"
                            )
                            with st.expander(student_label, expanded=False):
                                st.write(f"**Student Name:** {result.get('student_name', '')}")
                                st.write(f"**Student ID:** {result.get('student_id_parsed', '')}")
                                st.write(f"**Score:** {result.get('score', 0)} / {result.get('max_points', 0)}")
                                st.write(f"**Feedback:** {result.get('feedback', '')}")
                                st.write(f"**Model:** {result.get('model_name', '')}")
                                st.write(f"**Provider:** {result.get('model_provider', '')}")
                                st.markdown("**Detailed Explanation:**")
                                st.write(result.get("detailed_explanation", ""))

                    st.markdown("---")

                    # All session actions on one row: Load Setup on the left,
                    # Download Results in the middle, Delete on the far right.
                    # The CSV export is built unconditionally so the download
                    # button always has data ready even when session_results is
                    # empty (the button is simply not rendered in that case).
                    act_col1, act_col2, act_col3 = st.columns([2, 2, 1])

                    with act_col1:
                        # Load Setup — restores the exam setup (questions,
                        # rubric, sub-rubric, and max points) from this session
                        # into session state so the teacher can go directly to
                        # the Student Submissions tab and upload new student
                        # files without re-entering the exam setup.
                        # The input method is switched to "Enter questions manually"
                        # so the loaded questions text is immediately visible in
                        # the Setup Exam tab after the rerun.
                        if st.button("📂 Load Setup", key=f"eg_hist_load_{session_id}"):
                            setup = get_exam_grading_session_setup(session_id, user_id)
                            if setup:
                                st.session_state.questions              = setup.get("questions_text", "")
                                st.session_state.rubric                 = setup.get("rubric", "")
                                st.session_state.sub_rubric             = setup.get("sub_rubric", "")
                                st.session_state.max_points             = setup.get("max_points", 100)
                                st.session_state.questions_input_method = "Enter questions manually"  # show loaded text immediately
                                st.toast(
                                    "Exam setup loaded. Switch to the Setup Exam tab "
                                    "to review, then go to Student Submissions to grade."
                                )
                                st.rerun()
                            else:
                                st.error("Could not load setup data for this session.")

                    with act_col2:
                        # Build the CSV export from session results stored in
                        # the database. Column names and structure match the CSV
                        # produced by the Grading Results tab exactly so that
                        # downloads from history are interchangeable with live exports.
                        if session_results:
                            export_data = []
                            for r in session_results:
                                score   = r.get("score", 0)
                                max_pts = r.get("max_points", 0)
                                pct = (
                                    float(score) / float(max_pts) * 100
                                    if max_pts and float(max_pts) > 0
                                    else 0.0
                                )
                                export_data.append({
                                    "Student Name":      r.get("student_name", ""),
                                    "Student ID":        r.get("student_id_parsed", ""),
                                    "Score":             score,
                                    "Max Points":        max_pts,
                                    "Percentage":        pct,
                                    "Feedback":          r.get("feedback", "N/A"),
                                    "Detailed Feedback": r.get("detailed_explanation", "N/A"),
                                })
                            csv_data = pd.DataFrame(export_data).to_csv(index=False)
                            st.download_button(
                                label="📥 Download Results as CSV",
                                data=csv_data,
                                file_name=f"grading_results_{session_id[:8]}.csv",
                                mime="text/csv",
                                key=f"eg_hist_dl_{session_id}",
                            )

                    with act_col3:
                        if st.button(
                            "🗑️ Delete",
                            key=f"eg_hist_del_{session_id}",
                            type="primary",
                        ):
                            _dialog_delete_grading_session(session_id, user_id)