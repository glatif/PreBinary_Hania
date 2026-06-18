import streamlit as st
import PyPDF2
import docx
import pptx
import io

# --- Helper Functions for Text Extraction ---
def extract_text_from_pdf(file):
    """Extracts text from a PDF file."""
    try:
        pdf_reader = PyPDF2.PdfReader(file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() or ""
        return text
    except Exception as e:
        return f"Error reading PDF: {e}"

def extract_text_from_docx(file):
    """Extracts text from a DOCX file."""
    try:
        doc = docx.Document(file)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        return f"Error reading DOCX: {e}"

def extract_text_from_pptx(file):
    """Extracts text from a PPTX file."""
    try:
        presentation = pptx.Presentation(file)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error reading PPTX: {e}"

def extract_text_from_txt(file):
    """Extracts text from a TXT file."""
    try:
        return file.getvalue().decode("utf-8")
    except Exception as e:
        return f"Error reading TXT: {e}"

# --- Dummy Data ---
dummy_questions = [
    {
        "question": "What is the powerhouse of the cell?",
        "options": ["Mitochondria", "Nucleus", "Ribosome", "Chloroplast"],
        "correct_answer": "Mitochondria"
    },
    {
        "question": "Which planet is known as the Red Planet?",
        "options": ["Earth", "Mars", "Jupiter", "Saturn"],
        "correct_answer": "Mars"
    }
]

# --- Streamlit App UI ---
st.set_page_config(page_title="U-REAP Practice Questions", page_icon="🧠", layout="wide")

st.title("U-REAP: Your Personal AI Practice Question Generator")
st.markdown("Upload your study materials, and I'll create practice questions to help you master the content.")

# --- Initialize session state for analysis status ---
if 'analysis_complete' not in st.session_state:
    st.session_state.analysis_complete = False

# --- Tabbed Interface ---
tab1, tab2 = st.tabs(["📝 Setup Quiz", "🧠 Take Quiz"])

with tab1:
    st.header("🗂️ Step 1: Upload & Configure")
    
    # --- File Upload ---
    uploaded_files = st.file_uploader(
        "Upload PDFs, Word documents, presentations, or text files.",
        type=["pdf", "docx", "pptx", "txt"],
        accept_multiple_files=True,
        key="file_uploader"
    )

    # --- Customization ---
    st.header("🧠 Customization")
    st.markdown("Tailor your practice session to your needs. _(For display only in this version)_.")
    
    col1, col2 = st.columns(2)
    with col1:
        question_type = st.selectbox("Choose question type", ["Multiple Choice", "True/False", "Short Answer"], key="question_type")
        num_questions = st.number_input("Number of questions", 1, 50, 5, key="num_questions")
    with col2:
        difficulty = st.select_slider("Difficulty level", options=["Easy", "Medium", "Hard"], value="Medium", key="difficulty")
        topic_filters = st.text_input("Filter by topics (comma-separated)", placeholder="e.g., Biology, History", key="topic_filters")

    # --- Analysis ---
    analyze_button = st.button("Analyze Content & Generate Quiz", type="primary", key="analyze_button")
    if analyze_button and uploaded_files:
        st.session_state.full_context = ""
        st.session_state.analysis_complete = False
        with st.spinner('Analyzing documents... Please wait.'):
            all_texts = []
            for file in uploaded_files:
                file_extension = file.name.split('.')[-1].lower()
                text = ""
                if file_extension == "pdf":
                    text = extract_text_from_pdf(file)
                elif file_extension == "docx":
                    text = extract_text_from_docx(file)
                elif file_extension == "pptx":
                    text = extract_text_from_pptx(file)
                elif file_extension == "txt":
                    text = extract_text_from_txt(file)
                
                if text:
                    all_texts.append((file.name, text))

            if all_texts:
                st.session_state.full_context = "\n\n---\n\n".join([t[1] for t in all_texts])
                st.session_state.analysis_complete = True
                st.session_state.extracted_texts = all_texts
                st.success("Analysis complete! Head to the 'Take Quiz' tab.")
            else:
                st.error("Could not extract any text from the uploaded files.")

    # Display extracted text if analysis is complete
    if st.session_state.analysis_complete:
        with st.expander("View Extracted Content from Files", expanded=False):
            for file_name, text in st.session_state.get('extracted_texts', []):
                with st.container():
                    st.subheader(f"📄 {file_name}")
                    st.text_area("Content", text, height=200, key=f"text_{file_name}")

with tab2:
    st.header("🎯 Step 2: Test Your Knowledge")
    if not st.session_state.analysis_complete:
        st.warning("Please upload and analyze your documents in the 'Setup Quiz' tab first.")
    else:
        st.markdown("Here are some practice questions based on your materials. _(Currently using dummy questions)_.")
        for i, q in enumerate(dummy_questions):
            st.subheader(f"Question {i + 1}: {q['question']}")
            user_answer = st.radio(
                "Choose your answer:",
                options=q['options'],
                key=f"q_{i}",
                index=None # No default selection
            )

            if st.button(f"Check Answer for Q{i+1}", key=f"check_{i}"):
                if user_answer == q['correct_answer']:
                    st.success(f"Correct! The answer is {q['correct_answer']}.")
                elif user_answer is None:
                    st.warning("Please select an answer.")
                else:
                    st.error(f"Incorrect. The correct answer is {q['correct_answer']}.")
            st.markdown("---")
